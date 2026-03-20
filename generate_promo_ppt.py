from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/data/proact-vl")
OUT = ROOT / "Proact-VL_promo_safe.pptx"

BLUE = RGBColor(37, 99, 235)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(90, 104, 117)
LIGHT = RGBColor(243, 244, 246)
WHITE = RGBColor(255, 255, 255)
ORANGE_BG = RGBColor(253, 243, 231)
ORANGE_LINE = RGBColor(244, 198, 141)


def set_run(run, size=20, bold=False, color=TEXT, font="Arial"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_text(slide, left, top, width, height, text, size=20, bold=False,
             color=TEXT, align=PP_ALIGN.LEFT, font="Arial"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def add_multiline(slide, left, top, width, height, lines):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    first = True
    for text, size, bold, color in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        set_run(run, size=size, bold=bold, color=color)
        first = False
    return box


def add_box(slide, left, top, width, height, fill, line):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    return shape


def add_title(slide, title, subtitle=""):
    add_text(slide, Inches(0.65), Inches(0.35), Inches(11.5), Inches(0.5), title, size=28, bold=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(0.95), Inches(1.15), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    if subtitle:
        add_text(slide, Inches(0.65), Inches(1.05), Inches(11.5), Inches(0.35), subtitle, size=12, color=MUTED)


def add_bullets(slide, left, top, width, height, title, bullets, fill=WHITE, line=RGBColor(229, 231, 235)):
    add_box(slide, left, top, width, height, fill, line)
    box = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.16), width - Inches(0.44), height - Inches(0.32))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run(run, size=22, bold=True, color=BLUE)
    for item in bullets:
        p = tf.add_paragraph()
        p.text = item
        p.bullet = True
        p.level = 0
        if p.runs:
            set_run(p.runs[0], size=18, color=TEXT)
    return box


def fit_image(path, max_w, max_h):
    img = Image.open(path)
    w, h = img.size
    ratio = min(max_w / w, max_h / h)
    return Inches(w * ratio / 96.0), Inches(h * ratio / 96.0)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# 1 cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_multiline(
    slide, Inches(0.75), Inches(0.7), Inches(6.6), Inches(1.8),
    [
        ("Proact-VL", 30, True, BLUE),
        ("A Proactive VideoLLM for Real-Time AI Companions", 24, True, TEXT),
    ],
)
add_text(
    slide, Inches(0.78), Inches(2.05), Inches(8.6), Inches(0.6),
    "Weicai Yan*, Yuhong Dai*, Qi Ran, Haodong Li, Wang Lin, Hao Liao, Xing Xie, Tao Jin, Jianxun Lian",
    size=16, color=MUTED
)
add_text(slide, Inches(0.78), Inches(2.45), Inches(2.5), Inches(0.3), "* Equal contribution", size=11, color=MUTED)
add_box(slide, Inches(0.78), Inches(2.95), Inches(5.65), Inches(1.15), ORANGE_BG, ORANGE_LINE)
add_text(slide, Inches(1.02), Inches(3.2), Inches(0.4), Inches(0.3), "*", size=24, bold=True, color=BLUE)
add_text(
    slide, Inches(1.35), Inches(3.12), Inches(4.8), Inches(0.7),
    "Proact-VL turns multimodal language models into proactive, real-time AI companions that continuously watch, understand, decide when to speak, and generate timely commentary.",
    size=18, color=TEXT
)
pic_w, pic_h = fit_image(ROOT / "overview.png", 4.9 * 96, 5.2 * 96)
slide.shapes.add_picture(str(ROOT / "overview.png"), Inches(8.0), Inches(0.7), width=pic_w, height=pic_h)


# 2 why
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Why Proact-VL", "From passive video QA to proactive AI companionship")
add_bullets(
    slide, Inches(0.7), Inches(1.55), Inches(5.95), Inches(4.65), "Current gap",
    [
        "Most VideoLLMs are reactive and wait for prompts.",
        "Streaming environments require low latency over unbounded input.",
        "Companions must decide when to speak, not only what to say.",
    ]
)
add_bullets(
    slide, Inches(6.75), Inches(1.55), Inches(5.9), Inches(4.65), "What Proact-VL enables",
    [
        "Continuous perception over infinite video streams.",
        "Proactive response timing and triggering.",
        "Commentary, co-commentary, and guidance interaction.",
        "Flexible support for multiple Qwen-VL backbones.",
    ],
    fill=ORANGE_BG, line=ORANGE_LINE
)


# 3 key features
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Key Features")
cards = [
    ("Real-Time Processing", "Handles infinite video streams with low latency for continuous interaction."),
    ("Multi-Modal Commentary", "Supports single-speaker, multi-speaker, and guidance commentary scenarios."),
    ("Proactive Understanding", "Goes beyond reactive responses with contextual and timely intervention."),
    ("Flexible Architecture", "Works with Qwen2-VL, Qwen2.5-VL, and Qwen3-VL backbones."),
]
coords = [(0.7, 1.55), (6.7, 1.55), (0.7, 4.0), (6.7, 4.0)]
for (title, desc), (x, y) in zip(cards, coords):
    add_box(slide, Inches(x), Inches(y), Inches(5.9), Inches(1.9), WHITE, RGBColor(229, 231, 235))
    add_text(slide, Inches(x + 0.25), Inches(y + 0.2), Inches(5.3), Inches(0.35), title, size=20, bold=True)
    add_text(slide, Inches(x + 0.25), Inches(y + 0.62), Inches(5.25), Inches(0.85), desc, size=15, color=MUTED)


# 4 framework
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Framework Overview", "Streaming perception + response triggering + commentary generation")
pic_w, pic_h = fit_image(ROOT / "framework.png", 11.7 * 96, 4.7 * 96)
slide.shapes.add_picture(str(ROOT / "framework.png"), Inches((13.333 - (pic_w.inches)) / 2), Inches(1.45), width=pic_w, height=pic_h)
add_text(
    slide, Inches(0.9), Inches(6.35), Inches(11.6), Inches(0.5),
    "Proact-VL jointly addresses continuous perception, proactive response decision-making, and controllable real-time generation.",
    size=16, color=MUTED, align=PP_ALIGN.CENTER
)


# 5 scenarios
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Representative Scenarios")
titles = ["Live Commentary", "Co-Commentary", "Interactive Guidance"]
descs = [
    "A single agent watches a live stream and delivers timely commentary.",
    "Multiple speakers interact around the same stream with richer timing needs.",
    "The system proactively helps users with context-aware guidance.",
]
for i in range(3):
    x = 0.7 + i * 4.15
    add_box(slide, Inches(x), Inches(1.65), Inches(3.6), Inches(4.2), LIGHT, RGBColor(229, 231, 235))
    add_text(slide, Inches(x + 0.2), Inches(1.88), Inches(3.1), Inches(0.35), titles[i], size=20, bold=True, color=BLUE)
    add_text(slide, Inches(x + 0.2), Inches(2.3), Inches(3.0), Inches(0.9), descs[i], size=14, color=TEXT)
    add_box(slide, Inches(x + 0.28), Inches(3.2), Inches(3.0), Inches(1.7), WHITE, RGBColor(203, 213, 225))
    add_text(slide, Inches(x + 0.5), Inches(3.8), Inches(2.6), Inches(0.45), "Insert screenshot or demo frame", size=14, color=MUTED, align=PP_ALIGN.CENTER)


# 6 benchmark
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Evaluation Message")
add_bullets(
    slide, Inches(0.7), Inches(1.55), Inches(5.9), Inches(4.6), "Live Gaming Benchmark",
    [
        "Realistic settings for commentary and guidance.",
        "Measures both response quality and timing.",
        "Uses LLM-based judging for practical interaction evaluation.",
    ],
    fill=ORANGE_BG, line=ORANGE_LINE
)
add_bullets(
    slide, Inches(6.75), Inches(1.55), Inches(5.9), Inches(4.6), "Pitch takeaway",
    [
        "Not just another VideoLLM benchmark entry.",
        "Reframes the target from passive understanding to proactive companionship.",
        "Useful for both research-facing and product-facing storytelling.",
    ]
)


# 7 demo
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Demo Placeholder")
add_box(slide, Inches(1.0), Inches(1.65), Inches(11.3), Inches(4.9), LIGHT, RGBColor(203, 213, 225))
add_text(slide, Inches(3.9), Inches(3.1), Inches(5.0), Inches(0.5), "Demo video placeholder", size=26, bold=True, align=PP_ALIGN.CENTER)
add_text(
    slide, Inches(2.2), Inches(3.75), Inches(8.6), Inches(0.6),
    "Replace with a local video screenshot strip, exported GIF frames, or a linked showcase video.",
    size=18, color=MUTED, align=PP_ALIGN.CENTER
)


# 8 closing
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_multiline(
    slide, Inches(0.9), Inches(1.05), Inches(11.2), Inches(2.0),
    [
        ("Proact-VL", 30, True, BLUE),
        ("Toward proactive, real-time multimodal AI companions", 24, True, TEXT),
    ],
)
add_text(slide, Inches(0.92), Inches(3.0), Inches(9.0), Inches(0.4), "Paper: arXiv 2603.03447", size=20, color=MUTED)
add_text(slide, Inches(0.92), Inches(3.45), Inches(9.5), Inches(0.4), "Code: github.com/microsoft/AnthropomorphicIntelligence", size=20, color=MUTED)
add_text(slide, Inches(0.92), Inches(5.5), Inches(10.8), Inches(0.5), "Promo deck baseline. Replace the demo placeholder with your final assets when ready.", size=16, color=MUTED)


prs.save(str(OUT))
print(OUT)
